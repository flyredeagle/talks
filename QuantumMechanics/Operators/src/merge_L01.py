"""
merge_L01.py
Appends all slides from L01_extra_slides.pptx into ModuleI1_L01.pptx,
preserving the original deck's layout/design, then saves the combined
deck as ModuleI1_L01_extended.pptx
"""
import copy
from lxml import etree
from pptx import Presentation
from pptx.util import Pt
import os

SRC_MAIN  = "/home/claude/qm-module-i1/generated/ModuleI1_L01.pptx"
SRC_EXTRA = "/home/claude/L01_extra_slides.pptx"
OUT_PATH  = "/home/claude/ModuleI1_L01_extended.pptx"

def merge_presentations(base_path, extra_path, out_path):
    base  = Presentation(base_path)
    extra = Presentation(extra_path)

    # The slide layout XML is self-contained in each slide.
    # The simplest reliable approach: copy the XML element of each
    # extra slide, add a blank slide to base, then replace its spTree.

    slide_width  = base.slide_width
    slide_height = base.slide_height

    print(f"Base:  {len(base.slides)} slides  ({base_path})")
    print(f"Extra: {len(extra.slides)} slides ({extra_path})")

    for i, extra_slide in enumerate(extra.slides):
        # Pick a blank-ish layout from base (index 6 = blank in most themes)
        layout_idx = min(6, len(base.slide_layouts) - 1)
        new_slide = base.slides.add_slide(base.slide_layouts[layout_idx])

        # Replace the shape tree with the extra slide's shape tree
        sp_tree = new_slide.shapes._spTree
        # Remove all children except the first two (required namespace nodes)
        for child in list(sp_tree):
            sp_tree.remove(child)

        # Deep-copy every child from the extra slide's spTree
        extra_sp_tree = extra_slide.shapes._spTree
        for child in extra_sp_tree:
            sp_tree.append(copy.deepcopy(child))

        # Copy background if the extra slide has one set explicitly
        try:
            extra_bg = extra_slide.background
            extra_fill = extra_bg.fill
            if extra_fill.type is not None:
                new_bg = new_slide.background
                new_fill = new_bg.fill
                # Copy the background XML element
                bg_xml = extra_slide._element.find('.//{http://schemas.openxmlformats.org/drawingml/2006/spreadsheetDrawing}bg')
                # Simpler: copy the cSld/bg element
                ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
                extra_csld = extra_slide._element.find(f'{{{ns}}}cSld')
                new_csld   = new_slide._element.find(f'{{{ns}}}cSld')
                if extra_csld is not None and new_csld is not None:
                    extra_bg_elem = extra_csld.find(f'{{{ns}}}bg')
                    if extra_bg_elem is not None:
                        old_bg = new_csld.find(f'{{{ns}}}bg')
                        if old_bg is not None:
                            new_csld.remove(old_bg)
                        new_csld.insert(0, copy.deepcopy(extra_bg_elem))
        except Exception:
            pass  # background copy is best-effort

        if (i+1) % 5 == 0:
            print(f"  Copied slide {i+1}/{len(extra.slides)}...")

    print(f"\nTotal slides in merged deck: {len(base.slides)}")
    base.save(out_path)
    print(f"✓ Saved: {out_path}")

if __name__ == "__main__":
    merge_presentations(SRC_MAIN, SRC_EXTRA, OUT_PATH)

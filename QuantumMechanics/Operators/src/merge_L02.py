"""
merge_L02.py  —  Append L02 extra slides into the base L02 deck.
Input:  generated/ModuleI1_L02.pptx  (14 slides)
        L02_extra_slides.pptx         (29 slides)
Output: generated/ModuleI1_L02_extended.pptx (43 slides)
"""
import copy, os
from pptx import Presentation

REPO = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
SRC_MAIN  = os.path.join(REPO, "generated", "ModuleI1_L02.pptx")
SRC_EXTRA = os.path.join(REPO, "..", "L02_extra_slides.pptx")
OUT_PATH  = os.path.join(REPO, "generated", "ModuleI1_L02_extended.pptx")

def merge(base_path, extra_path, out_path):
    base  = Presentation(base_path)
    extra = Presentation(extra_path)
    print(f"Base:  {len(base.slides)} slides  ({base_path})")
    print(f"Extra: {len(extra.slides)} slides ({extra_path})")
    for i, extra_slide in enumerate(extra.slides):
        layout_idx = min(6, len(base.slide_layouts)-1)
        new_slide  = base.slides.add_slide(base.slide_layouts[layout_idx])
        sp_tree    = new_slide.shapes._spTree
        for child in list(sp_tree):
            sp_tree.remove(child)
        for child in extra_slide.shapes._spTree:
            sp_tree.append(copy.deepcopy(child))
        try:
            ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
            ec = extra_slide._element.find(f'{{{ns}}}cSld')
            nc = new_slide._element.find(f'{{{ns}}}cSld')
            if ec is not None and nc is not None:
                eb = ec.find(f'{{{ns}}}bg')
                if eb is not None:
                    ob = nc.find(f'{{{ns}}}bg')
                    if ob is not None:
                        nc.remove(ob)
                    nc.insert(0, copy.deepcopy(eb))
        except Exception:
            pass
        if (i+1) % 5 == 0:
            print(f"  Copied {i+1}/{len(extra.slides)} slides...")
    print(f"\nTotal slides: {len(base.slides)}")
    base.save(out_path)
    print(f"✓ Saved: {out_path}")

if __name__ == "__main__":
    merge(SRC_MAIN, SRC_EXTRA, OUT_PATH)
